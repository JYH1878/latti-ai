from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1a, 0x1a, 0x1a)
    p.alignment = PP_ALIGN.CENTER
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.3), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(prs, title, bullets, table_data=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(0x2c, 0x3e, 0x50)
    title_shape.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
    left = Inches(0.5)
    top = Inches(1.4)
    width = Inches(12.3)
    height = Inches(5.8)
    if table_data:
        width = Inches(6.0)
    content_box = slide.shapes.add_textbox(left, top, width, height)
    tf = content_box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(14)
        p.level = 0
    if table_data:
        rows = len(table_data)
        cols = len(table_data[0])
        table = slide.shapes.add_table(rows, cols, Inches(6.7), Inches(1.4), Inches(6.0), Inches(0.8 + 0.5*rows)).table
        for i, row in enumerate(table_data):
            for j, cell_text in enumerate(row):
                cell = table.cell(i, j)
                cell.text = str(cell_text)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    if i == 0:
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                    else:
                        paragraph.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0x2c, 0x3e, 0x50)
    return slide

add_title_slide(prs, "任务5：稀疏打包 Bootstrapping 算子优化", "LattiAI / LattiSense 开源框架中的可配置稀疏密钥封装\n\n CipherFlow 盗火者计划 · 赛道三")

add_content_slide(prs, "背景与动机", [
    "Bootstrapping 占 FHE 推理时间的 60%~90%，是最大瓶颈",
    "标准 Dense 密钥（H=N/2）的 CtS/StC 乘法深度高、旋转密钥庞大",
    "稀疏打包（SSE）：ModRaise 阶段临时切换到低汉明权重密钥（H=32）",
    "LattiAI 已集成 Lattigo，但缺少：可配置化、系统性量化对比、编译器自适应",
    "目标：让稀疏打包 Bootstrapping 可配置、可对比、可自适应"
])

add_content_slide(prs, "技术方案：三层递进优化", [
    "1. SDK 参数可配置化 — Go/C++ SDK 支持 8 种预设参数切换（低风险）",
    "2. BSGS 调优验证 — Sparse vs Dense 全 ratio 对比 + 生产级验证（中风险）",
    "3. 编译器放置优化 — score 修正 + 子图阈值动态调整（低风险）"
])

add_content_slide(prs, "Sparse Secret Encapsulation 原理", [
    "1. ModRaise：提升模数链至 Q_L",
    "2. KeySwitch Dense→Sparse：切换到低权重密钥 sk_sparse（H=32）",
    "3. CoeffsToSlots (CtS)：稀疏密钥下执行同态 DFT，乘法深度↓",
    "4. EvalMod：同态模约简（Sine/Cosine 逼近）",
    "5. SlotsToCoeffs (StC)：稀疏密钥下执行逆 DFT",
    "6. KeySwitch Sparse→Dense：切换回标准密度密钥",
    "收益：旋转次数、密钥交换、模数消耗均显著下降"
])

add_content_slide(prs, "代码改动点", [
    "Go SDK (bootstrap.go)：8 种 preset 枚举 + CreateCkksBtpParameterByPreset",
    "C++ SDK (fhe_lib_v2.h/.cpp)：create_parameter_by_preset(int)",
    "编译器 (components.py)：新增 N16QP1547H192H32（32.1 bits 高精度）",
    "编译器 (pipeline.py)：btp_param_list 扩展为多参数自动选优",
    "编译器 (score.py)：Sparse bootstrap 折扣 + 动态子图阈值",
    "Benchmark：sparse_bench_test.go / bsgs_opt_test.go / sparse_vs_dense_bsgs_test.go"
])

add_content_slide(prs, "实验环境", [
    "CPU：AMD Ryzen 9 7945HX · Go 1.24 · Lattigo v4 backend",
    "测试方法：Go benchmark，benchtime=1x，分阶段计时",
    "Sparse：N16QP1546H192H32（H=192/H=32）",
    "Dense：N16QP1767H32768H32（H=N/2/H=32）"
])

add_content_slide(prs, "实验结果：Sparse vs Dense 性能对比", [
    "Sparse SSE 实现 2.85 倍加速，节省 4 个模数层级"
], table_data=[
    ["阶段", "Dense", "Sparse", "加速比"],
    ["ModUp", "1.45 s", "0.30 s", "4.8×"],
    ["CtS", "37.66 s", "11.93 s", "3.2×"],
    ["EvalMod", "9.47 s", "5.84 s", "1.6×"],
    ["StC", "12.43 s", "3.37 s", "3.7×"],
    ["总时间", "~61.0 s", "~21.4 s", "2.85×"],
    ["模数层级", "29", "25", "节省 4 层"]
])

add_content_slide(prs, "实验结果：精度与安全性", [
    "精度：N16QP1546H192H32 → 26.6 bits；N16QP1547H192H32 → 32.1 bits",
    "安全性：128-bit 安全级别，H=192 满足 RLWE 要求",
    "临时密钥 H=32 仅用于 bootstrapping 内部，通过 KeySwitch 隔离",
    "失败概率：2^-138.7（可忽略）—— Boura et al. 2022 严格证明"
])

add_content_slide(prs, "BSGS Ratio 调优与 Sparse/Dense 对比验证", [
    "Sparse 在所有 ratio 下均优于 Dense",
    "ratio 越小，Sparse 优势越大（1.26× → 1.13×）",
    "在完整生产级参数（LogN=16）下，CtS+StC 实现 2.04× 加速",
    "旋转密钥数相同（70个），差异完全来自稀疏密钥封装的线性变换效率提升"
], table_data=[
    ["Ratio", "Sparse", "Dense", "加速比"],
    ["1.0", "685.7 ms", "863.6 ms", "1.26×"],
    ["2.0", "742.2 ms", "898.4 ms", "1.21×"],
    ["4.0", "791.0 ms", "897.2 ms", "1.13×"]
])

add_content_slide(prs, "编译器层 Bootstrap 放置优化（双层）", [
    "改进1——时间估算修正（score.py）：",
    "  · 新增 is_sparse_bootstrapping_param() + SPARSE_BTP_DISCOUNT = 0.40",
    "  · BtpScoreParam.get_score() 自动应用折扣",
    "改进2——子图阈值动态调整（graph_partition_dp.py）：",
    "  · Sparse 参数下 level_threshold 从 max_depth-4 放宽至 max_depth-6",
    "  · 允许探索更多子图候选，找到更优放置方案",
    "收益：编译器决策更贴近实际性能"
])

add_content_slide(prs, "编译器层收益：打破单一参数锁定", [
    "改动前：btp_param_list = [N16QP1546H192H32]，所有模型锁定",
    "改动后：btp_param_list = [N16QP1546H192H32, N16QP1547H192H32]，自动选优",
    "场景1（追求速度）：N16QP1546H192H32，残差层级多（420 bits）",
    "场景2（追求精度）：N16QP1547H192H32，默认 scale 更高（2^45）"
])

add_content_slide(prs, "成果总结", [
    "可配置化：Go/C++ SDK 新增 8 种预设参数接口",
    "可量化：Sparse vs Dense 端到端 benchmark，2.85× 加速",
    "可扩展：编译器层支持多参数候选 + 放置策略优化",
    "可调优：BSGS ratio 调优实验验证 Sparse 在全 ratio 下优于 Dense",
    "可优化：编译器放置策略引入 Sparse 折扣，决策更贴近实际性能"
])

add_content_slide(prs, "展望", [
    "LCR + AKS（Kim et al. 2025）：预计再节省 1 个模数层级，加速 1.28×",
    "自适应参数选择：根据模型深度、slot 数、精度阈值自动选择最优预设",
    "GPU 后端适配：在 HEonGPU 中针对稀疏密钥封装做 kernel 特化",
    "BSGS 自动调优：根据矩阵稀疏度自动搜索最优 BSGS ratio"
])

prs.save('/home/jyh/latti-ai/任务5_稀疏打包Bootstrapping优化.pptx')
print("PPT generated successfully")
